import re
import json
import collections
import collections.abc
import docx2txt
from io import StringIO

from pptx import Presentation
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfdocument import PDFDocument
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.pdfpage import PDFPage
from pdfminer.pdfparser import PDFParser


class Search:
    """This class provides methods that search strings for PII based on file type."""

    def __init__(self):
        pass

    def char_search(self, a_string):
        """This function checks a string for social security numbers and credit card numbers.

        Args:
            a_string (str): the string to search in

        Returns:
            int: returns 0 if no PII found
            tuple: key,value pair if PII found
        """
        pii_values = {
            "Amex": "^3[47][0-9]{13}$",
            "BCGlobal": "^(6541|6556)[0-9]{12}$",
            "Carte Blanche": "^389[0-9]{11}$",
            "Diners Club": "^3(?:0[0-5]|[68][0-9])[0-9]{11}$",
            "Discover": "^65[4-9][0-9]{13}|64[4-9][0-9]{13}|6011[0-9]{12}|(622(?:12[6-9]|1[3-9][0-9]|[2-8][0-9][0-9]|9[01][0-9]|92[0-5])[0-9]{10})$",
            "Insta Payment": "^63[7-9][0-9]{13}$",
            "JCB": "^(?:2131|1800|35\d{3})\d{11}$",
            "Laser": "^(6304|6706|6709|6771)[0-9]{12,15}$",
            "Maestro": "^(5018|5020|5038|6304|6759|6761|6763)[0-9]{8,15}$",
            "Master": "^(5[1-5][0-9]{14}|2(22[1-9][0-9]{12}|2[3-9][0-9]{13}|[3-6][0-9]{14}|7[0-1][0-9]{13}|720[0-9]{12}))$",
            "Solo": "^(6334|6767)[0-9]{12}|(6334|6767)[0-9]{14}|(6334|6767)[0-9]{15}$",
            "Switch": "^(4903|4905|4911|4936|6333|6759)[0-9]{12}|(4903|4905|4911|4936|6333|6759)[0-9]{14}|(4903|4905|4911|4936|6333|6759)[0-9]{15}|564182[0-9]{10}|564182[0-9]{12}|564182[0-9]{13}|633110[0-9]{10}|633110[0-9]{12}|633110[0-9]{13}$",
            "Union Pay": "^(62[0-9]{14,17})$",
            "Visa": "^4[0-9]{12}(?:[0-9]{3})?$",
            "Bank Routing Number": "\b((0[0-9])|(1[0-2])|(2[1-9])|(3[0-2])|(6[1-9])|(7[0-2])|80)([0-9]{7})\b",
            "Swift Code": "\b[A-Z]{6}[A-Z0-9]{2}([A-Z0-9]{3})?\b",
            "Social Security Number": "^(?!000|666)[0-8][0-9]{2}-(?!00)[0-9]{2}-(?!0000)[0-9]{4}$",
            "ITIN": "^9\d{2}-?((5[0-9]|6[0-5])|(8[3-8])|(9[0-2])|(9[4-9]))-?\d{4}$",
            "MRN Type 1": "[1-9]{3}-[1-9]{1}-[1-9]{5}",
        }
        a_string_list = a_string.split()
        for section in a_string_list:
            # Here is where you check if it matches the parameters
            for key, value in pii_values.items():
                value_search = re.search(str(value), section)
                if value_search:
                    return key, section
        return 0

    def line_cleanup_and_check(self, the_list, line_number, filename, type_check):
        if type_check == "txt":
            occurrences = open("exposed_txt.txt", "a")
        if type_check == "json":
            occurrences = open("exposed_json.txt", "a")
        if type_check == "pdf":
            occurrences = open("exposed_pdf.txt", "a")
        if type_check == "csv":
            occurrences = open("exposed_csv.txt", "a")
        if type_check == "docx":
            occurrences = open("exposed_docx.txt", "a")
        if type_check == "pptx":
            occurrences = open("exposed_pptx.txt", "a")

        second_count = 0
        if len(the_list) > 0:
            for section in the_list:
                check_letters = section.isalpha()
                second_count += 1
                if (not check_letters) and (len(section) > 7) and (len(section) < 17):
                    check_pii_value = self.char_search(section)
                    if check_pii_value != 0:
                        if line_number == 0:
                            occurrences.write(
                                f"File: {filename} Value:{check_pii_value} Line: {second_count}\n"
                            )
                        else:
                            occurrences.write(
                                f"File: {filename} Value: {check_pii_value} Line: {line_number}\n"
                            )
                        occurrences.close()
                        return
        return

    def txt_search(self, file_read):
        """Searches for PII in a txt file

        Args:
            file_read (FILE): A white-space delimited txt file
        """

        # Uncomment following lines if you want to see what files
        # were checked
        # txt_file = open("txt_files_checked.txt", "a")
        # txt_file.write(file_read.strip("./") + "\n")

        with open(file_read, "r") as file:
            linecount = 0
            line = file.readline()
            for line in file:
                linecount += 1
                line_list = [x.strip() for x in line.split()]
                self.line_cleanup_and_check(line_list, linecount, file_read, "txt")

        file.close()
        # txt_file.close()

    def csv_search(self, file_read):
        """Searches for PII in a CSV file.

        Args:
            file_read (CSV): A comma delimited file.
        """
        # Uncomment following lines if you want to see what files
        # were checked
        # csv_file = open("csv_files_checked.txt", "a")
        # csv_file.write(file_read.strip("./") + "\n")

        with open(file_read, "r") as file:

            linecount = 0
            line = file.readline()
            for line in file:
                linecount += 1
                line_list = [x.strip() for x in line.split(",")]
                self.line_cleanup_and_check(line_list, linecount, file_read, "csv")

        file.close()
        # csv_file.close()

    def convert_pdf_to_string(self, file_read):
        """Converts text from a PDF file into a string.

        Args:
            file_read (FILE): A PDF file

        Returns:
            str: The text from the PDF as a string
        """
        output_string = StringIO()

        with open(file_read, "rb") as in_file:
            parser = PDFParser(in_file)
            doc = PDFDocument(parser)
            rsrcmgr = PDFResourceManager()
            device = TextConverter(rsrcmgr, output_string, laparams=LAParams())
            interpreter = PDFPageInterpreter(rsrcmgr, device)
            for page in PDFPage.create_pages(doc):
                interpreter.process_page(page)
        in_file.close()
        return output_string.getvalue()

    def pdf_search(self, file_read):
        """Searches for PII in a PDF file.

        Args:
            file_read (PDF FILE): A PDF file
        """
        # uncomment following lines if you want to see what files were checked
        # don't forget to close file

        # pdf_file = open("pdf_files_checked.txt", "a")
        # pdf_file.write(filename[0] + "\n")

        filename = file_read.split(".//")
        text = self.convert_pdf_to_string(file_read)
        text_list = text.split()
        self.line_cleanup_and_check(text_list, 0, filename, "pdf")

        # pdf_file.close()

    def docx_search(self, file_read):
        my_text = docx2txt.process(file_read)
        # uncomment following lines if you want to see what files were checked
        # don't forget to close file

        # docx_file = open("docx_files_checked.txt", "a")
        # docx_file.write(file_read.strip("./") + "\n")

        line_list = [x.strip() for x in my_text.split()]
        self.line_cleanup_and_check(line_list, 0, file_read, "docx")

        # docx_file.close()

    def pptx_search(self, file_read):
        file = Presentation(file_read)
        for slide in file.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    line = shape.text
                    line_list = line.split()
                    self.line_cleanup_and_check(line_list, 0, file_read, "pptx")
        file.close()

    def json_search(self, file_read):
        """Searches for PII in a JSON file.

        Args:
            file_read (JSON FILE): A JSON file
        """
        result = None
        occurrences = open("exposed_json.txt", "a", encoding="utf-8")
        filename = file_read.split(".//")
        with open(file_read) as data_file:
            data = json.load(data_file)

            for value in data.values():
                if isinstance(value, str):
                    if (len(value) > 7) and (len(value) < 17):
                        result = self.char_search(value)

                if isinstance(value, int):
                    new_string = str(value)
                    if (len(new_string) > 7) and (len(new_string) < 17):
                        result = self.char_search(new_string)

                if isinstance(value, list):
                    for item in value:
                        string_to_search = str(item)
                        if (len(string_to_search) > 7) and (len(string_to_search) < 17):
                            result = self.char_search(string_to_search)

                if isinstance(value, dict):
                    for key, val in value.items():
                        new_string = str(val)
                        if (len(new_string) > 7) and (len(new_string) < 17):
                            result = self.char_search(new_string)

                if (result != 0) and (result != None):
                    occurrences.write(f"File: {filename[0]} Value: {result}\n")
                    occurrences.close()
                    return
        data_file.close()
        occurrences.close()
